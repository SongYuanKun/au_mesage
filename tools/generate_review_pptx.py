from __future__ import annotations

import datetime as dt
import pathlib

import yaml
from pptx import Presentation


ROOT = pathlib.Path(__file__).resolve().parents[1]
REQ_PATH = ROOT / "docs" / "SRS" / "requirements.yml"
OUT_DIR = ROOT / "docs" / "SRS" / "attachments" / "review"


def _add_title(prs: Presentation, title: str, subtitle: str) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = title
    slide.placeholders[1].text = subtitle


def _add_bullets(prs: Presentation, title: str, bullets: list[str]) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = title
    tf = slide.shapes.placeholders[1].text_frame
    tf.clear()
    for i, b in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = b


def main() -> None:
    data = yaml.safe_load(REQ_PATH.read_text(encoding="utf-8"))
    meta = data["meta"]
    roles = data.get("roles", [])
    fr = data.get("functional_requirements", [])
    nfr = data.get("non_functional_requirements", [])

    now = dt.datetime.now(dt.timezone.utc).astimezone().strftime("%Y-%m-%d")

    prs = Presentation()
    _add_title(
        prs,
        f"{meta.get('product')}｜需求评审报告",
        f"版本 {meta.get('version')} · {now} · 状态 {meta.get('status')}",
    )

    _add_bullets(
        prs,
        "背景与目标",
        [
            "解决贵金属价格信息分散、口径不一致与缺少提醒闭环的问题",
            "以“可信数据 + 实时提醒 + 任务型工具”形成差异化",
            "本次评审聚焦：需求可测试、无二义性、100%可追溯",
        ],
    )

    _add_bullets(
        prs,
        "角色与范围",
        [
            "角色：" + "、".join(r["name"] for r in roles),
            f"功能性需求条数：{len(fr)}；非功能性需求条数：{len(nfr)}",
            "覆盖模块：用户端、管理端、后台服务、接口与安全合规",
        ],
    )

    must_fr = [r for r in fr if str(r.get("moscow")).lower() == "must"]
    should_fr = [r for r in fr if str(r.get("moscow")).lower() == "should"]
    _add_bullets(
        prs,
        "优先级概览（MoSCoW）",
        [
            f"Must（本期不可缺）：{len(must_fr)} 条",
            f"Should（重要可后置）：{len(should_fr)} 条",
            "建议先交付：数据可信度展示 + SSE提醒稳定性 + 核心概览/趋势",
        ],
    )

    top_risks = []
    for r in fr + nfr:
        risks = r.get("risks") or []
        if risks:
            top_risks.append(f"{r['id']}：{risks[0]}")
        if len(top_risks) >= 6:
            break
    if not top_risks:
        top_risks = ["暂无"]
    _add_bullets(prs, "主要风险（Top）", top_risks)

    _add_bullets(
        prs,
        "评审结论与行动项",
        [
            "确认需求清单与MoSCoW分级",
            "确认接口字段口径与错误码规范",
            "确认RTM字段：DesignRef/CodeRef/TestRef落地方式",
            "确认运维与可观测指标的SLO/告警阈值",
        ],
    )

    OUT_DIR.mkdir(parents=True, exist_ok=True)
    out_path = OUT_DIR / "需求评审报告.pptx"
    prs.save(out_path)


if __name__ == "__main__":
    main()

